"""
Document processor for handling various file formats
"""
import logging
from pathlib import Path
from typing import List, Dict, Any
import pypdf
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process various document formats into text chunks"""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def process_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from PDF files"""
        try:
            chunks = []
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                full_text = ""

                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    full_text += text + "\n"

                # Split into chunks
                text_chunks = self._create_chunks(full_text)

                for idx, chunk in enumerate(text_chunks):
                    chunks.append({
                        "text": chunk,
                        "source": str(file_path),
                        "source_type": "pdf",
                        "chunk_id": idx,
                        "metadata": {
                            "total_pages": len(pdf_reader.pages),
                            "filename": file_path.name
                        }
                    })

            logger.info(f"Processed PDF {file_path.name}: {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []

    def process_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from DOCX files"""
        try:
            chunks = []
            doc = Document(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs])

            text_chunks = self._create_chunks(full_text)

            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "text": chunk,
                    "source": str(file_path),
                    "source_type": "docx",
                    "chunk_id": idx,
                    "metadata": {
                        "filename": file_path.name
                    }
                })

            logger.info(f"Processed DOCX {file_path.name}: {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return []

    def process_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint files"""
        try:
            chunks = []
            prs = Presentation(file_path)
            full_text = ""

            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                full_text += slide_text

            text_chunks = self._create_chunks(full_text)

            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "text": chunk,
                    "source": str(file_path),
                    "source_type": "pptx",
                    "chunk_id": idx,
                    "metadata": {
                        "total_slides": len(prs.slides),
                        "filename": file_path.name
                    }
                })

            logger.info(f"Processed PPTX {file_path.name}: {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return []

    def process_txt(self, file_path: Path) -> List[Dict[str, Any]]:
        """Process plain text files"""
        try:
            chunks = []
            with open(file_path, 'r', encoding='utf-8') as file:
                full_text = file.read()

            text_chunks = self._create_chunks(full_text)

            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "text": chunk,
                    "source": str(file_path),
                    "source_type": "txt",
                    "chunk_id": idx,
                    "metadata": {
                        "filename": file_path.name
                    }
                })

            logger.info(f"Processed TXT {file_path.name}: {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error processing TXT {file_path}: {e}")
            return []

    def process_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Process a file based on its extension"""
        suffix = file_path.suffix.lower()

        processors = {
            '.pdf': self.process_pdf,
            '.docx': self.process_docx,
            '.pptx': self.process_pptx,
            '.txt': self.process_txt,
            '.md': self.process_txt,
        }

        processor = processors.get(suffix)
        if processor:
            return processor(file_path)
        else:
            logger.warning(f"Unsupported file format: {suffix}")
            return []

    def _create_chunks(self, text: str) -> List[str]:
        """Split text into overlapping chunks"""
        if not text or len(text.strip()) == 0:
            return []

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + self.chunk_size

            # Try to break at sentence boundary
            if end < text_length:
                # Look for period, question mark, or exclamation point
                for i in range(end, max(start, end - 100), -1):
                    if text[i] in '.!?\n':
                        end = i + 1
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - self.chunk_overlap

        return chunks
